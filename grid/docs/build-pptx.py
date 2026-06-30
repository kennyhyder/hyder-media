#!/usr/bin/env python3
"""Build GridScout investor slide deck as PowerPoint with embedded videos."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

PURPLE = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
DARK_BG = RGBColor(0x0F, 0x0A, 0x1A)
DEEP_BG = RGBColor(0x1E, 0x10, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

MP4_DIR = os.path.join(os.path.dirname(__file__), "mp4")

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# Use blank layout
blank_layout = prs.slide_layouts[6]


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, left, top, width, height, font_size=14, bold=False,
             color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=13, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '+'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': '7C3AED'})
        buClr.append(srgb)
        buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Calibri'})
        buSz = pPr.makeelement(qn('a:buSzPct'), {'val': '120000'})
        pPr.append(buSz)
        pPr.append(buFont)
        pPr.append(buClr)
        pPr.append(buChar)
        # indent
        pPr.set('marL', str(Emu(Inches(0.3))))
        pPr.set('indent', str(Emu(Inches(-0.25))))
    return txBox


def add_badge(slide, text, left, top, bg_color, text_color):
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Calibri"
    # Badge bg via shape fill
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = bg_color
    # Round the corners if possible -- not directly supported, skip


POSTER_DIR = os.path.join(os.path.dirname(__file__), "posters")

def add_video(slide, video_path, left, top, width, height):
    """Add video with a real poster frame extracted from the video."""
    video_name = os.path.splitext(os.path.basename(video_path))[0]
    poster_path = os.path.join(POSTER_DIR, f"{video_name}.jpg")
    if not os.path.exists(poster_path):
        # fallback to a blank image
        from PIL import Image
        poster_path = "/tmp/gridscout_poster.png"
        img = Image.new('RGB', (1280, 720), (15, 10, 26))
        img.save(poster_path)

    shape = slide.shapes.add_movie(
        video_path, left, top, width, height,
        poster_frame_image=poster_path,
        mime_type='video/mp4'
    )
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text(slide, "\u26A1", Inches(0), Inches(1.8), W, Inches(0.8),
         font_size=52, alignment=PP_ALIGN.CENTER)
add_text(slide, "GridScout", Inches(0), Inches(2.6), W, Inches(1),
         font_size=52, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "Data Center Site Intelligence Platform",
         Inches(0), Inches(3.5), W, Inches(0.6),
         font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text(slide, "154,000+ SCORED SITES  \u2022  50 STATES  \u2022  14 SCORING DIMENSIONS",
         Inches(0), Inches(4.5), W, Inches(0.4),
         font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE DEFINITIONS
# ============================================================
slides_data = [
    {
        "badge": "DASHBOARD",
        "badge_color": PURPLE,
        "title": "Executive Overview",
        "desc": "A single-screen summary of the entire opportunity landscape. Score distributions, geographic concentration, and the top-ranked sites across all 50 states.",
        "bullets": [
            "154K+ sites scored across substations, industrial, and greenfield categories",
            "Score distribution histogram with average, median, min, and max statistics",
            "Top 15 states ranked by average DC readiness score",
            "Top 25 highest-scoring sites with power, fiber, and capacity metrics",
            "8 authoritative federal data sources cited",
        ],
        "video": "Screen Recording-Dashboard.mp4",
        "video_side": "left",
        "bg": DARK_BG,
    },
    {
        "badge": "INTERACTIVE MAP",
        "badge_color": PURPLE,
        "title": "Visual Infrastructure Analysis",
        "desc": "A full-screen geospatial explorer with 7 independent data layers. Instantly correlate site opportunities with power, fiber, and existing data center infrastructure.",
        "bullets": [
            "7 toggleable layers: sites, DCs, IXPs, transmission, substations, fiber, ISO queues",
            "Heat map mode reveals macro patterns of infrastructure suitability",
            "Location search with radius-based proximity filtering",
            "Real-time score distribution and site type breakdown",
            "3 base maps: dark, satellite, and street view",
            "Viewport-based dynamic loading for 150K+ markers",
        ],
        "video": "Screen Recording-Map.mp4",
        "video_side": "right",
        "bg": DEEP_BG,
    },
    {
        "badge": "GREENFIELD SITES",
        "badge_color": GREEN,
        "title": "Site Search & Screening",
        "desc": "The primary screening tool for narrowing thousands of candidate sites to a qualified shortlist. Customize the scoring model to match your investment thesis.",
        "bullets": [
            "12 sortable columns: score, voltage, capacity, IXP distance, energy price, flood zone",
            "Adjustable weight editor to re-score sites against your own criteria in real time",
            "Compare mode for side-by-side evaluation of up to 10 sites",
            "Saved shortlists for investment committee workflows",
            "ISO region filter: PJM, MISO, ERCOT, CAISO, SPP, ISO-NE, NYISO",
            "CSV export for offline analysis",
        ],
        "video": "Screen Recording-Greenfields.mp4",
        "video_side": "left",
        "bg": DARK_BG,
    },
    {
        "badge": "INDUSTRIAL SITES",
        "badge_color": AMBER,
        "title": "Brownfield Opportunities",
        "desc": "Decommissioned power plants with existing grid connections, cleared land, and road access. These sites can shave 2\u20134 years off the timeline to energization.",
        "bullets": [
            "Retired coal, gas, and nuclear plants with known MW capacity",
            "Cleanup status tracking: completed, in-progress, or pending",
            "Nearest substation distance for interconnection feasibility",
            "Acreage, retirement date, and former-use classification",
            "Map overlay showing all industrial sites with popup details",
            "Sortable by capacity, acreage, substation proximity",
        ],
        "video": "Screen Recording-Industrial.mp4",
        "video_side": "right",
        "bg": DEEP_BG,
    },
    {
        "badge": "TRANSMISSION LINES",
        "badge_color": PURPLE,
        "title": "Power Delivery Infrastructure",
        "desc": "Comprehensive view of the US transmission network. Identify high-voltage routes, upgrade candidates, and utility ownership for interconnection planning.",
        "bullets": [
            "Voltage class filtering: 100kV through 500kV+",
            "Upgrade candidate identification (50\u2013100 MW lines eligible for reconductoring)",
            "Utility ownership lookup for interconnection planning",
            "Interactive map rendering 2,500+ lines with voltage-weighted thickness",
            "Route geometry from substation to substation",
            "Status tracking: in-service, proposed, under construction",
        ],
        "video": "Screen Recording-Transmission.mp4",
        "video_side": "left",
        "bg": DARK_BG,
    },
    {
        "badge": "ENERGY CORRIDORS",
        "badge_color": GREEN,
        "title": "Pre-Permitted Federal Land",
        "desc": "Federal energy corridors with completed environmental review. Section 368 corridors, NIETC designations, and BLM DLAs dramatically reduce permitting timelines.",
        "bullets": [
            "Section 368 pre-approved corridors across 12 western states",
            "NIETC areas with FERC backstop siting authority",
            "BLM Solar DLAs pre-screened for energy development on federal land",
            "Acreage, capacity, and transmission line counts per corridor",
            "Environmental review status for each corridor",
        ],
        "video": "Screen Recording-Corridors.mp4",
        "video_side": "right",
        "bg": DEEP_BG,
    },
    {
        "badge": "HYPERSCALE TRACKER",
        "badge_color": AMBER,
        "title": "Competitive Landscape",
        "desc": "Tracks frontier AI and cloud data center construction by every major hyperscaler. Understand where demand is concentrating and where grid congestion risks are emerging.",
        "bullets": [
            "13 operators: Microsoft, Google, Meta, Amazon, xAI, CoreWeave, Oracle, and more",
            "Pipeline status: operational, under construction, and announced projects",
            "Planned capacity in GW per operator and per state",
            "Clickable operator and status breakdowns filter the project table",
            "Linked to GridScout site search for each state with active deployments",
        ],
        "video": "Screen Recording-Hyperscale.mp4",
        "video_side": "left",
        "bg": DARK_BG,
    },
    {
        "badge": "SITE DETAIL",
        "badge_color": GREEN,
        "title": "Investment-Grade Dossier",
        "desc": "A complete site evaluation for any of the 154K+ scored locations. Every data point sourced from federal agencies with deep links for verification.",
        "bullets": [
            "14-factor score breakdown with expandable data sources and scoring criteria",
            "Auto-generated investment thesis: strengths and risk flags",
            "Due diligence checklist: verified, flagged, and manual verification items",
            "Interactive map with transmission lines and fiber routes rendered",
            "Nearby IXPs and datacenters with contacts and distances",
            "Speed-to-energization analysis with ISO queue depth and wait times",
            "Land acquisition guidance with owner contacts and listing links",
            "50+ deep links to FEMA, EPA, EIA, WRI, FCC, PeeringDB, and more",
        ],
        "video": "Screen Recording-site-detail.mp4",
        "video_side": "right",
        "bg": DEEP_BG,
    },
]

total_slides = len(slides_data) + 2  # +title +closing

for i, sd in enumerate(slides_data):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, sd["bg"])

    slide_num = i + 2
    # Slide number
    add_text(slide, f"{slide_num} / {total_slides}",
             Inches(12.2), Inches(7.0), Inches(1), Inches(0.3),
             font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

    # Layout dimensions
    margin = Inches(0.7)
    video_w = Inches(7.0)
    video_h = Inches(3.94)  # 16:9 aspect
    text_w = Inches(5.0)

    if sd["video_side"] == "left":
        vid_left = margin
        txt_left = margin + video_w + Inches(0.4)
    else:
        txt_left = margin
        vid_left = margin + text_w + Inches(0.4)

    vid_top = Inches(1.6)
    txt_top = Inches(1.2)

    # Video
    video_path = os.path.join(MP4_DIR, sd["video"])
    if os.path.exists(video_path):
        add_video(slide, video_path, vid_left, vid_top, video_w, video_h)

    # Badge
    bc = sd["badge_color"]
    badge_bg = RGBColor(
        min(255, bc[0] // 3),
        min(255, bc[1] // 3),
        min(255, bc[2] // 3),
    )
    add_badge(slide, sd["badge"], txt_left, txt_top, badge_bg, sd["badge_color"])

    # Title
    add_text(slide, sd["title"], txt_left, txt_top + Inches(0.45), text_w, Inches(0.6),
             font_size=28, bold=True, color=WHITE)

    # Description
    add_text(slide, sd["desc"], txt_left, txt_top + Inches(1.1), text_w, Inches(1.0),
             font_size=12, color=LIGHT_GRAY)

    # Bullets
    bullet_top = txt_top + Inches(2.0)
    add_bullet_list(slide, sd["bullets"], txt_left, bullet_top, text_w, Inches(3.5),
                    font_size=11, color=LIGHT_GRAY)


# ============================================================
# SLIDE 10: Closing
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, DARK_BG)

add_text(slide, "Built for the Opportunity",
         Inches(0), Inches(1.5), W, Inches(0.8),
         font_size=40, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, "GridScout turns 154,000 potential sites into a ranked, filterable, investment-ready pipeline.",
         Inches(2), Inches(2.3), Inches(9.333), Inches(0.6),
         font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Three value prop cards
card_w = Inches(3.6)
card_h = Inches(2.2)
card_top = Inches(3.5)
gap = Inches(0.4)
start_left = (W - card_w * 3 - gap * 2) / 2

cards = [
    ("14 Scoring Dimensions", "Power, fiber, water, hazard, labor, climate, tax incentives, and more. Every score is auditable back to its federal data source."),
    ("8 Federal Data Sources", "HIFLD, FEMA NRI, WRI Aqueduct, BLS, FCC, EIA, NOAA, and PeeringDB. Updated and cross-referenced continuously."),
    ("Weeks to Minutes", "Auto-generated investment thesis, due diligence checklist, and land acquisition guidance for every site in the database."),
]

for j, (title, desc) in enumerate(cards):
    left = start_left + j * (card_w + gap)
    # Card background
    shape = slide.shapes.add_shape(
        1, int(left), card_top, card_w, card_h  # MSO_SHAPE.RECTANGLE = 1
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x12, 0x2A)
    shape.line.fill.background()

    add_text(slide, title, int(left) + Inches(0.3), card_top + Inches(0.3),
             card_w - Inches(0.6), Inches(0.4),
             font_size=16, bold=True, color=LIGHT_PURPLE)
    add_text(slide, desc, int(left) + Inches(0.3), card_top + Inches(0.8),
             card_w - Inches(0.6), Inches(1.2),
             font_size=11, color=LIGHT_GRAY)

add_text(slide, "KENNY HYDER  \u2022  KENNY@HYDER.ME  \u2022  HYDER.ME/GRID",
         Inches(0), Inches(6.2), W, Inches(0.4),
         font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text(slide, f"{total_slides} / {total_slides}",
         Inches(12.2), Inches(7.0), Inches(1), Inches(0.3),
         font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

# Save
out_path = os.path.join(os.path.dirname(__file__), "GridScout-Platform-Overview.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Size: {os.path.getsize(out_path) / 1024 / 1024:.1f} MB")
